#!/usr/bin/env python3
"""Generate FastGA Deep Dive & Vibe Optimization presentation using template."""

import os
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Paths ──────────────────────────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DOCS_DIR = os.path.dirname(SCRIPT_DIR)
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "2025-10-22-triton_seq_align.pptx")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "fastga_deep_dive.pptx")

IMG_STORAGE_HUMAN = os.path.join(DOCS_DIR, "storage_timeline_human_T32.png")
IMG_OPT1 = os.path.join(DOCS_DIR, "storage_optimization",
                         "storage_timeline_early_gix_comparison.png")
IMG_OPT3 = os.path.join(DOCS_DIR, "storage_optimization",
                         "storage_timeline_opt3_comparison.png")

# ── Colors ─────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# ── Template Layout Indices ────────────────────────────────────────────────
# Layout 0: "Title Slide"     - placeholders: 0=title, 1=subtitle
# Layout 1: "Title and Content" - placeholders: 0=title, 1=content, 12=slide_number
# Layout 2: "Section Header"  - placeholders: 0=title, 1=text
# Layout 3: "Title Only"      - placeholder: 0=title
LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_TITLE_ONLY = 3

# ── Dimensions (from template: 13.33 x 7.50 inches) ───────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.67)  # match template content placeholder left edge
CONTENT_TOP = Inches(1.09)  # match template content placeholder top
CONTENT_W = Inches(10.97)   # match template content width
CONTENT_H = Inches(4.95)    # match template content height


# ── Helpers ────────────────────────────────────────────────────────────────

def delete_all_slides(prs):
    """Remove all existing slides from the presentation."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_run(run, text, size=18, bold=False, italic=False, color=None):
    """Set text properties on a run. Color=None means inherit from theme."""
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color


def fill_placeholder(slide, idx, text, size=None):
    """Fill a placeholder by index with text."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            if size:
                run.font.size = Pt(size)
            return ph
    return None


def add_content_bullets(slide, bullets, sub_bullets=None, size=18):
    """Fill the content placeholder (idx=1) with bullet points."""
    sub_bullets = sub_bullets or {}
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                p.level = 0
                set_run(p.add_run(), bullet, size=size)

                if i in sub_bullets:
                    for sub in sub_bullets[i]:
                        sp = tf.add_paragraph()
                        sp.space_after = Pt(4)
                        sp.level = 1
                        set_run(sp.add_run(), sub, size=size - 2)
            return tf
    return None


def add_image_centered(slide, img_path, top=None, max_width=None, max_height=None):
    """Add an image centered horizontally, scaled to fit."""
    top = top or CONTENT_TOP
    max_width = max_width or CONTENT_W
    max_height = max_height or (SLIDE_H - top - Inches(0.5))

    from PIL import Image as PILImage
    with PILImage.open(img_path) as im:
        img_w, img_h = im.size
    aspect = img_w / img_h

    if max_width / aspect <= max_height:
        w = max_width
        h = int(w / aspect)
    else:
        h = max_height
        w = int(h * aspect)

    left = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, left, top, w, h)


def add_table(slide, headers, rows, left=None, top=None, width=None,
              row_height=Inches(0.4)):
    """Add a formatted table."""
    left = left or MARGIN
    top = top or CONTENT_TOP
    width = width or CONTENT_W
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = row_height * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    col_width = width // n_cols
    for i in range(n_cols):
        table.columns[i].width = col_width

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p.add_run(), h, size=14, bold=True, color=WHITE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 1 else WHITE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            text = str(val)
            is_bold = text.startswith("**") and text.endswith("**")
            if is_bold:
                text = text[2:-2]
            set_run(p.add_run(), text, size=13, bold=is_bold, color=CHARCOAL)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ── Slide Builders ─────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Title slide using template Layout 0."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    fill_placeholder(slide, 0, "FastGA Deep Dive\n& Vibe Optimization")
    fill_placeholder(slide, 1, "Reducing Storage Footprint for Whole-Genome Alignment on HPC\n\nJiajie Li  |  March 2026")
    return slide


def slide_section(prs, title, subtitle=""):
    """Section divider using template Layout 2."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    fill_placeholder(slide, 0, subtitle)
    fill_placeholder(slide, 1, title, size=36)
    return slide


def slide_content(prs, title, bullets, sub_bullets=None, size=18):
    """Standard content slide using template Layout 1."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    fill_placeholder(slide, 0, title)
    add_content_bullets(slide, bullets, sub_bullets=sub_bullets, size=size)
    return slide


def slide_content_with_table(prs, title):
    """Content slide with title only (table added separately)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE_ONLY])
    fill_placeholder(slide, 0, title)
    return slide


def slide_03_problem(prs):
    return slide_content(prs, "Whole-Genome Alignment Is Hard", [
        "Goal: find all local alignments between two genome assemblies",
        "Scale: human genomes have ~3 billion base pairs each",
        "Existing tools (minimap2, MUMmer) use hash-table lookups for seed finding",
        "Hash tables cause random memory access \u2192 CPU cache misses",
        "At genome scale, the CPU spends most time waiting for RAM (~200 cycles/miss)",
    ])


def slide_04_key_idea(prs):
    return slide_content(prs, "FastGA: Cache Coherence Over Hash Tables", [
        "By Gene Myers, Richard Durbin, and Chenxi Zhou",
        "Core principle: replace hash-table lookups with sorted arrays + linear merges",
        "Pre-sort all k-mers (K=40) using MSD radix sort (not comparison-based)",
        "Find seeds by merging two sorted indices in a single linear sweep",
        "Every memory access is sequential and predictable \u2192 CPU cache stays warm",
        "Compact output: .1aln trace-point format is 15\u201316\u00d7 smaller than PAF+CIGAR",
    ])


def slide_06_pipeline(prs):
    return slide_content(prs, "4-Stage Pipeline", [
        "Stage 1: Compress \u2014 FASTA \u2192 GDB (2-bit packed DNA)",
        "Stage 2: Index \u2014 GDB \u2192 GIX (sorted k-mer index, syncmer-filtered)",
        "Stage 3: Seed Merge \u2014 linear merge of two sorted GIX indices",
        "Stage 4: Sort, Chain, Align \u2014 radix sort seeds \u2192 chain \u2192 wave-front alignment",
    ], sub_bullets={
        0: ["\u2022 4 bases per byte. ~750 MB for a human genome (vs ~3 GB FASTA)"],
        1: ["\u2022 ~10 GB/Gbp. 15 bytes/entry: 7B suffix + 1B mask + 1B LCP + 4B pos + 2B contig"],
        2: ["\u2022 Longest prefix match + frequency filter (\u2264 10 occurrences)"],
        3: ["\u2022 O(nd) wave-front algorithm within diagonal band defined by chain"],
    })


def slide_06b_subphases(prs):
    slide = slide_content_with_table(prs, "Pipeline Sub-Phases & Output Files")
    add_table(slide,
              ["Sub-phase", "Tool", "What it does", "Output files", "Threading"],
              [
                  ["Compress", "FAtoGDB (\u00d72)", "FASTA \u2192 2-bit packed binary", ".1gdb + .bps", "Single"],
                  ["Index", "GIXmake (\u00d72)", "Syncmer k-mers \u2192 MSD radix sort", ".gix + .ktab.*", "Multi (max 32)"],
                  ["Seed Merge", "FastGA ph.1", "Linear merge of two sorted GIX", "_pair.*.N/.C (temp)", "Multi"],
                  ["Sort + Align", "FastGA ph.2\u20134", "Radix sort \u2192 chain \u2192 align \u2192 merge", ".1aln (final output)", "Multi"],
                  ["Convert", "ALNtoPAF", "Decode trace points \u2192 PAF text", "PAF to stdout", "Single"],
              ],
              top=Inches(1.5), row_height=Inches(0.55))

    # File lifecycle table below
    add_table(slide,
              ["File", "Size (human, per genome)", "Created by", "Read by", "Deleted"],
              [
                  [".1gdb + .bps", "~750 MB", "FAtoGDB", "FastGA (base fetching)", "At exit (-k keeps)"],
                  [".gix + .ktab.*", "**~32 GB**", "GIXmake", "FastGA (seed merge only)", "At exit"],
                  ["_pair.* (temp)", "~3.5 GB", "Seed merge", "Sort+align", "Unlinked after open"],
                  [".1aln", "157 MB", "Sort+align", "ALNtoPAF", "Persistent output"],
              ],
              top=Inches(4.5), row_height=Inches(0.45))
    return slide


def slide_07_data_structures(prs):
    slide = slide_content_with_table(prs, "Key Data Structures & File Sizes")
    add_table(slide,
              ["Format", "Contents", "Size Scaling", "Lifetime"],
              [
                  ["GDB (.1gdb + .bps)", "2-bit compressed DNA", "~0.25 GB/Gbp", "Entire run"],
                  ["GIX (.gix + .ktab.*)", "Sorted k-mer index + LCP", "**~10 GB/Gbp**", "Seed merge only"],
                  ["_pair.* (temp)", "Seed position pairs", "~2.3 GB/Gbp", "Sort+align only"],
                  [".1aln (output)", "Trace-point alignments", "~0.05 GB/Gbp", "Persistent"],
              ],
              top=Inches(1.6))

    tb = add_textbox(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "GIX dominates storage: 62.7 GB for two human genomes (88% of peak disk). "
            "Temp files are invisible to ls (unlinked immediately after open).",
            size=15, italic=True, color=CHARCOAL)
    return slide


def slide_09_performance(prs):
    slide = slide_content_with_table(prs, "Human Genome Performance (GRCh38 vs CHM13, T=32)")

    add_table(slide,
              ["Metric", "Value"],
              [
                  ["Wall clock", "**9 min 48s**"],
                  ["Peak RSS", "19.0 GB"],
                  ["Total seeds found", "1.3 billion"],
                  ["Non-redundant alignments", "518,037 (avg 28.4 Kbp)"],
                  ["CPU utilization", "638% (of 3200%)"],
              ],
              left=MARGIN, top=Inches(1.6), width=Inches(5.5))

    add_table(slide,
              ["Phase", "Wall Time", "% Total"],
              [
                  ["GDB creation", "18s", "3%"],
                  ["GIX build", "82s", "14%"],
                  ["Seed merge", "5.3s", "1%"],
                  ["**Sort + align**", "**482s (8 min)**", "**82%**"],
              ],
              left=Inches(7.0), top=Inches(1.6), width=Inches(5.7))

    tb = add_textbox(slide, MARGIN, Inches(5.3), CONTENT_W, Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "Sort+align dominates at human scale (82% of runtime). "
            "Seed merge is extremely fast (5.3s) thanks to the cache-coherent linear sweep.",
            size=15, italic=True, color=CHARCOAL)
    return slide


def slide_10_storage_problem(prs):
    return slide_content(prs, "The Storage Problem", [
        "Peak disk usage: 71 GB for one human-vs-human comparison",
        "GIX indices: 62.7 GB = 88% of peak storage",
        "GIX is only read during the 5.3-second seed merge phase",
        "But GIX stays on disk through the entire 8-minute sort+align phase",
        "Two concurrent runs on one HPC node: 142 GB of scratch needed",
        "Key insight: 88% of storage is held for 98% of runtime unnecessarily",
    ])


def slide_10b_minimap2_vs_fastga(prs):
    slide = slide_content_with_table(prs, "Design Tradeoff: minimap2 vs FastGA")
    add_table(slide,
              ["Aspect", "minimap2", "FastGA"],
              [
                  ["Index type", "Hash table (minimizers)", "Sorted k-mer table + LCP"],
                  ["Index size (human)", "~6\u20138 GB", "**~32 GB/genome**"],
                  ["Index location", "In memory", "On disk"],
                  ["Seed finding", "Hash lookup (random access)", "Linear merge (sequential)"],
                  ["Match type", "Fixed k-mer length only", "Variable-length (adaptamers)"],
                  ["Cache behavior", "Cache misses on every lookup", "Cache-friendly sequential scan"],
                  ["Reusable index", "No (rebuilt each run)", "Yes (-k flag)"],
                  ["Output format", "PAF (text)", ".1aln trace-points (15\u201316\u00d7 smaller)"],
              ],
              top=Inches(1.5), row_height=Inches(0.48))

    tb = add_textbox(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "FastGA trades ~4\u00d7 larger index for cache-coherent seed finding and "
            "free variable-length matching. The extra bytes (k-mer suffix + LCP) are "
            "the structural cost of the sorted-merge design.",
            size=15, italic=True, color=CHARCOAL)
    return slide


def slide_10c_gix_breakdown(prs):
    slide = slide_content_with_table(prs, "What\u2019s Inside the 63 GB GIX? (Human, Both Genomes)")
    add_table(slide,
              ["Field", "Bytes/entry", "Total (GB)", "% of GIX", "Purpose"],
              [
                  ["K-mer suffix", "7", "~32 GB", "**52%**", "Sorted DNA sequence for merge comparison"],
                  ["Position", "4", "~18 GB", "**29%**", "Where in the contig this k-mer starts"],
                  ["Contig+Strand", "2 or 1", "~7 GB", "11%", "Which contig, which strand"],
                  ["Mask byte", "1", "~5 GB", "7%", "Soft-masking (always 0 when off \u2014 waste)"],
                  ["LCP", "1", "~5 GB", "7%", "Longest common prefix with previous entry"],
                  ["Prefix tables", "fixed", "0.3 GB", "<1%", "Two 128 MB jump tables"],
              ],
              top=Inches(1.5), row_height=Inches(0.52))

    tb = add_textbox(slide, MARGIN, Inches(5.4), CONTENT_W, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    set_run(p.add_run(),
            "K-mer suffix + Position = 81% of GIX \u2014 fundamentally required by the algorithm. "
            "Optimization targets: mask byte (removed in Opt 3, \u22124.6 GB), "
            "LCP (could recompute on-the-fly), k-mer suffix (hard to compress \u2014 Opt 5 failed).",
            size=15, italic=True, color=CHARCOAL)
    return slide


def slide_11_storage_timeline(prs):
    slide = slide_content_with_table(prs, "Storage Timeline: Human Genomes (T=32)")
    if os.path.exists(IMG_STORAGE_HUMAN):
        add_image_centered(slide, IMG_STORAGE_HUMAN, top=Inches(1.4))
    return slide


def slide_13_opt_plan(prs):
    slide = slide_content_with_table(prs, "6 Identified Optimization Opportunities")
    add_table(slide,
              ["#", "Optimization", "Est. Impact (Human)", "Quality", "Status"],
              [
                  ["1", "Early GIX deletion", "Frees 63 GB in sort+align", "Bit-exact", "\u2705 Done"],
                  ["2", "Sparse prefix index", "~254 MB", "Bit-exact", "Not started"],
                  ["3", "Eliminate mask byte", "-4.6 GB (-7.7% ktab)", "Bit-exact", "\u2705 Done"],
                  ["4", "On-the-fly LCP", "~4.4 GB (~7% ktab)", "Bit-exact", "Deferred"],
                  ["5", "Ktab compression", "~14% ktab", "Bit-exact", "\u274c Failed"],
                  ["6", "Aggressive syncmer filter", "~40% ktab", "Trade-off", "Not started"],
              ],
              top=Inches(1.6))
    return slide


def slide_14_opt1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    fill_placeholder(slide, 0, "Opt 1: Early GIX Deletion (Zero Cost)")

    # Use content placeholder for bullets
    add_content_bullets(slide, [
        "Delete GIX immediately after seed merge, not at program exit",
        "Sort+align disk: 64 GB \u2192 7 GB (\u201389%)",
        "Two concurrent human runs: 128 GB \u2192 14 GB",
        "Zero performance impact, bit-exact output",
        "Only ~35 lines changed in FastGA.c",
    ])

    # Add image below content area
    if os.path.exists(IMG_OPT1):
        # Resize content placeholder to make room
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.height = Inches(2.8)
                break
        add_image_centered(slide, IMG_OPT1,
                           top=Inches(3.9), max_height=Inches(3.2))
    return slide


def slide_15_opt3_opt5(prs):
    slide = slide_content_with_table(prs, "Opt 3: Eliminate Mask Byte  |  Opt 5: Compression")

    # Opt 3 on the left
    tb1 = add_textbox(slide, MARGIN, Inches(1.6), Inches(5.5), Inches(3.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    set_run(p.add_run(), "Opt 3: Eliminate Mask Byte ", size=20, bold=True, color=NAVY)
    set_run(p.add_run(), "\u2705 Implemented", size=16, color=RGBColor(0x22, 0x8B, 0x22))
    for bullet in [
        "Remove unused 1-byte mask field from ktab entries (15B \u2192 14B)",
        "Result: 7.7% ktab reduction, bit-exact output",
        "Human genomes: \u22124.6 GB total GIX savings",
    ]:
        bp = tf1.add_paragraph()
        bp.space_after = Pt(6)
        set_run(bp.add_run(), bullet, size=17, color=CHARCOAL)

    # Opt 5 on the right
    tb2 = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(3.0))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    set_run(p2.add_run(), "Opt 5: Ktab Compression ", size=20, bold=True, color=NAVY)
    set_run(p2.add_run(), "\u274c Failed", size=16, color=RGBColor(0xCC, 0x00, 0x00))
    for bullet in [
        "Attempted zstd block compression of ktab partitions",
        "Only 14% compression (position data nearly random)",
        "Decompression round-trip bug: 53.6M vs 51.1M seeds",
        "All code reverted \u2014 adds dependency for modest savings",
    ]:
        bp = tf2.add_paragraph()
        bp.space_after = Pt(6)
        set_run(bp.add_run(), bullet, size=17, color=CHARCOAL)

    # Lesson learned
    tb3 = add_textbox(slide, MARGIN, Inches(4.8), CONTENT_W, Inches(1.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    set_run(tf3.paragraphs[0].add_run(),
            "Lesson: Not all optimizations work out. "
            "The position payload in ktab entries contains coordinate data that is "
            "effectively random \u2014 compression algorithms can't exploit patterns that don't exist.",
            size=15, italic=True, color=CHARCOAL)
    return slide


def slide_16_cumulative(prs):
    slide = slide_content_with_table(prs, "Cumulative Results: Opt 1 + Opt 3")

    add_table(slide,
              ["Metric", "Baseline", "Opt 1+3", "Savings"],
              [
                  ["GIX size (human, both)", "62.7 GB", "58.1 GB", "-4.6 GB (-7.3%)"],
                  ["Peak total disk (human)", "71 GB", "66.4 GB", "-4.6 GB (-6.5%)"],
                  ["Sort+align disk (human)", "**64 GB**", "**7 GB**", "**-57 GB (-89%)**"],
                  ["GIX size (EXAMPLE, both)", "1,864 MB", "1,740 MB", "-124 MB (-6.6%)"],
                  ["Sort+align disk (EXAMPLE)", "1,905 MB", "438 MB", "-1,467 MB (-77%)"],
              ],
              top=Inches(1.6), width=CONTENT_W)

    if os.path.exists(IMG_OPT3):
        add_image_centered(slide, IMG_OPT3,
                           top=Inches(4.3), max_height=Inches(2.8))
    return slide


def slide_17_reflections(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    fill_placeholder(slide, 0, "Reflections: What We Actually Did")

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()

            # "What worked" header
            p0 = tf.paragraphs[0]
            set_run(p0.add_run(), "What we did: removed low-hanging fruit", size=20, bold=True)

            for bullet in [
                "Opt 1: moved a delete() call earlier \u2014 trivial code change, huge practical impact",
                "Opt 2: changed KBYTES+1 to KBYTES when masking is off \u2014 one line of logic",
                "Both are bit-exact, zero performance cost, zero risk",
            ]:
                bp = tf.add_paragraph()
                bp.level = 0
                bp.space_after = Pt(6)
                set_run(bp.add_run(), bullet, size=17)

            # "Honest assessment" header
            hp = tf.add_paragraph()
            hp.space_before = Pt(16)
            set_run(hp.add_run(), "Honest assessment", size=20, bold=True)

            for bullet in [
                "These optimizations are trivial \u2014 we found waste and removed it",
                "We didn\u2019t make any fundamental algorithmic tradeoffs",
                "Our one non-trivial attempt (ktab compression) failed",
                "The real optimization space \u2014 trading sensitivity, compute, or quality for storage \u2014 remains unexplored",
            ]:
                bp = tf.add_paragraph()
                bp.level = 0
                bp.space_after = Pt(6)
                set_run(bp.add_run(), bullet, size=17)
            break
    return slide


def slide_18_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    fill_placeholder(slide, 0, "Next Steps: Exploring the Tradeoff Space")

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()

            p0 = tf.paragraphs[0]
            set_run(p0.add_run(), "Sensitivity \u2194 Storage", size=19, bold=True)
            for sub in [
                "Aggressive syncmer filtering: select ~40% instead of ~75% of positions \u2192 ~40% smaller GIX",
                "Tune frequency threshold (-f): higher = fewer seeds = less temp storage, but miss repeat-flanking matches",
            ]:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.space_after = Pt(4)
                set_run(sp.add_run(), sub, size=16)

            hp = tf.add_paragraph()
            hp.space_before = Pt(12)
            set_run(hp.add_run(), "Compute \u2194 Storage", size=19, bold=True)
            for sub in [
                "On-the-fly LCP: recompute during merge instead of storing (~7% GIX reduction, adds CPU cost)",
                "Streaming seed consumption: chain and align seeds as they\u2019re produced, avoid materializing all at once",
            ]:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.space_after = Pt(4)
                set_run(sp.add_run(), sub, size=16)

            hp2 = tf.add_paragraph()
            hp2.space_before = Pt(12)
            set_run(hp2.add_run(), "Toward autonomous optimization", size=19, bold=True)
            for sub in [
                "Define clear objectives: storage budget, sensitivity floor, runtime ceiling",
                "Let an AI agent explore the parameter/design space systematically",
                "Auto-benchmark, auto-evaluate, auto-iterate \u2014 \u00e0 la Karpathy\u2019s AutoResearch vision",
            ]:
                sp = tf.add_paragraph()
                sp.level = 1
                sp.space_after = Pt(4)
                set_run(sp.add_run(), sub, size=16)
            break
    return slide


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    # Load template and delete existing slides
    prs = Presentation(TEMPLATE_PATH)
    delete_all_slides(prs)

    # Section 1: Introduction
    slide_01_title(prs)
    slide_section(prs, "Introduction to FastGA")
    slide_03_problem(prs)
    slide_04_key_idea(prs)

    # Section 2: Pipeline
    slide_section(prs, "FastGA's Pipeline")
    slide_06_pipeline(prs)
    slide_06b_subphases(prs)
    slide_07_data_structures(prs)

    # Section 3: Benchmarks
    slide_section(prs, "Benchmark Results")
    slide_09_performance(prs)
    slide_10_storage_problem(prs)
    slide_10b_minimap2_vs_fastga(prs)
    slide_10c_gix_breakdown(prs)
    slide_11_storage_timeline(prs)

    # Section 4: Optimization
    slide_section(prs, "Storage Optimization Attempts")
    slide_13_opt_plan(prs)
    slide_14_opt1(prs)
    slide_15_opt3_opt5(prs)
    slide_16_cumulative(prs)

    # Conclusion
    slide_section(prs, "Reflections & Next Steps")
    slide_17_reflections(prs)
    slide_18_next_steps(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved to {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
